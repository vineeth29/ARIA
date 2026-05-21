import os, json, re, sys, subprocess, datetime

SCRIPT_DIR   = os.path.dirname(os.path.abspath(__file__))
CONFIG_FILE  = os.path.join(SCRIPT_DIR, "data", "docsearch_config.json")
INDEX_FILE   = os.path.join(SCRIPT_DIR, "data", "docsearch_index.json")
os.makedirs(os.path.join(SCRIPT_DIR, "data"), exist_ok=True)

SUPPORTED = {".pdf", ".pptx", ".ppt", ".docx", ".doc", ".txt", ".md",
             ".xlsx", ".xls", ".csv", ".png", ".jpg", ".jpeg", ".bmp", ".tiff"}

def _load_config():
    if os.path.exists(CONFIG_FILE):
        try:
            with open(CONFIG_FILE) as f:
                return json.load(f)
        except Exception:
            pass
    return {"folders": [], "indexed_at": None}

def _save_config(cfg):
    with open(CONFIG_FILE, "w") as f:
        json.dump(cfg, f, indent=2)

def _load_index():
    if os.path.exists(INDEX_FILE):
        try:
            with open(INDEX_FILE) as f:
                return json.load(f)
        except Exception:
            pass
    return {}

def _save_index(idx):
    with open(INDEX_FILE, "w") as f:
        json.dump(idx, f, indent=2)

def add_folder(path):
    path = os.path.expanduser(path.strip().strip('"').strip("'"))
    if not os.path.exists(path):
        return False, f"Folder not found: {path}"
    cfg = _load_config()
    if path not in cfg["folders"]:
        cfg["folders"].append(path)
        _save_config(cfg)
        return True, f"Added: {path}"
    return False, f"Already added: {path}"

def remove_folder(path):
    cfg = _load_config()
    if path in cfg["folders"]:
        cfg["folders"].remove(path)
        _save_config(cfg)
        return True
    return False

def get_folders():
    return _load_config().get("folders", [])

def list_folders():
    folders = get_folders()
    if not folders:
        return "No folders added yet.\nUse: /docs add C:\\path\\to\\folder"
    lines = ["Search folders:"]
    for i, f in enumerate(folders, 1):
        count = _count_files(f)
        lines.append(f"  {i}. {f}  ({count} supported files)")
    return "\n".join(lines)

def _count_files(folder):
    count = 0
    try:
        for root, dirs, files in os.walk(folder):
            dirs[:] = [d for d in dirs if not d.startswith('.')]
            for f in files:
                if os.path.splitext(f)[1].lower() in SUPPORTED:
                    count += 1
    except Exception:
        pass
    return count

def _extract_pdf(path):
    try:
        import fitz
        doc = fitz.open(path)
        pages = []
        for i, page in enumerate(doc):
            text = page.get_text().strip()
            if text and len(text) > 20:
                pages.append(text)
            else:
                pix = page.get_pixmap(dpi=150)
                img_path = path + f"_page{i}.png"
                pix.save(img_path)
                ocr = _ocr_image(img_path)
                try: os.unlink(img_path)
                except: pass
                if ocr:
                    pages.append(f"[Handwritten/Scanned Page {i+1}]: {ocr}")
        return "\n\n".join(pages)
    except ImportError:
        return _extract_pdf_fallback(path)
    except Exception as e:
        return None

def _extract_pdf_fallback(path):
    try:
        result = subprocess.run(
            ["pdftotext", path, "-"],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip()
    except Exception:
        pass
    return None

def _extract_pptx(path):
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
                if shape.shape_type == 13:
                    try:
                        from io import BytesIO
                        import tempfile
                        img_data = shape.image.blob
                        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
                        tmp.write(img_data)
                        tmp.close()
                        ocr = _ocr_image(tmp.name)
                        try: os.unlink(tmp.name)
                        except: pass
                        if ocr:
                            texts.append(f"[Image text]: {ocr}")
                    except Exception:
                        pass
            if texts:
                slides.append(f"[Slide {i}]: " + " | ".join(texts))
        return "\n".join(slides)
    except Exception:
        return None

def _extract_docx(path):
    try:
        from docx import Document
        doc = Document(path)
        parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n".join(parts)
    except Exception:
        return None

def _extract_image_ocr(path):
    return _ocr_image(path)

def _ocr_image(path):
    try:
        import pytesseract
        from PIL import Image
        img  = Image.open(path)
        text = pytesseract.image_to_string(img)
        return text.strip() if text.strip() else None
    except ImportError:
        pass
    try:
        result = subprocess.run(
            ["tesseract", path, "stdout", "--psm", "6"],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0:
            return result.stdout.strip()
    except Exception:
        pass
    return None

def _extract_csv(path):
    try:
        import csv
        rows = []
        with open(path, newline="", encoding="utf-8", errors="ignore") as f:
            for row in csv.reader(f):
                rows.append(" | ".join(row))
        return "\n".join(rows[:200])
    except Exception:
        return None

def _extract_text(path):
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read(50000)
    except Exception:
        return None

def extract_file_content(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return _extract_pdf(path)
    elif ext in (".pptx", ".ppt"):
        return _extract_pptx(path)
    elif ext in (".docx", ".doc"):
        return _extract_docx(path)
    elif ext in (".xlsx", ".xls"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            parts = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                parts.append(f"[Sheet: {sheet}]")
                for row in ws.iter_rows(values_only=True):
                    vals = [str(c) for c in row if c is not None]
                    if vals:
                        parts.append(" | ".join(vals))
            return "\n".join(parts)
        except Exception:
            return None
    elif ext == ".csv":
        return _extract_csv(path)
    elif ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff"):
        return _extract_image_ocr(path)
    else:
        return _extract_text(path)

def build_index(status_cb=None):
    folders = get_folders()
    if not folders:
        return 0, "No folders set. Use: /docs add <folder_path>"
    index = {}
    total = 0
    for folder in folders:
        if not os.path.exists(folder):
            continue
        for root, dirs, files in os.walk(folder):
            dirs[:] = [d for d in dirs if not d.startswith('.')]
            for fname in files:
                ext = os.path.splitext(fname)[1].lower()
                if ext not in SUPPORTED:
                    continue
                fpath = os.path.join(root, fname)
                try:
                    mtime = os.path.getmtime(fpath)
                    existing = index.get(fpath, {})
                    if existing.get("mtime") == mtime and existing.get("content"):
                        continue
                    if status_cb:
                        status_cb(f"Indexing: {fname}")
                    content = extract_file_content(fpath)
                    if content:
                        index[fpath] = {
                            "name":    fname,
                            "folder":  root,
                            "ext":     ext,
                            "mtime":   mtime,
                            "content": content[:20000],
                            "size":    len(content)
                        }
                        total += 1
                except Exception:
                    continue
    cfg = _load_config()
    cfg["indexed_at"] = datetime.datetime.now().isoformat()
    _save_config(cfg)
    _save_index(index)
    return total, f"Indexed {total} files from {len(folders)} folder(s)."

def _score(query_words, text):
    text_lower = text.lower()
    score = 0
    for word in query_words:
        count = text_lower.count(word)
        if count > 0:
            score += count
            score += 3
    return score

def search(query, top_n=5):
    index = _load_index()
    if not index:
        return [], "No index found. Type /docs index to build it first."
    stop = {"what","is","are","the","a","an","in","on","at","of","to","for",
            "and","or","tell","me","about","find","search","show","does","do",
            "which","where","how","when","who","my","i","can","you","please"}
    words = [w.lower() for w in re.findall(r'\b\w+\b', query) if w.lower() not in stop]
    if not words:
        return [], "Query too vague — try being more specific."
    results = []
    for fpath, data in index.items():
        content = data.get("content", "")
        score = _score(words, content + " " + data.get("name", ""))
        if score > 0:
            snippet = _get_snippet(content, words)
            results.append({
                "path":    fpath,
                "name":    data["name"],
                "ext":     data["ext"],
                "score":   score,
                "snippet": snippet,
                "content": content
            })
    results.sort(key=lambda x: x["score"], reverse=True)
    return results[:top_n], None

def _get_snippet(content, words, context=150):
    content_lower = content.lower()
    best_pos = -1
    for word in words:
        pos = content_lower.find(word)
        if pos != -1:
            best_pos = pos
            break
    if best_pos == -1:
        return content[:200].strip()
    start = max(0, best_pos - 60)
    end   = min(len(content), best_pos + context)
    snippet = content[start:end].strip()
    snippet = re.sub(r'\s+', ' ', snippet)
    if start > 0:
        snippet = "..." + snippet
    if end < len(content):
        snippet = snippet + "..."
    return snippet

def format_search_results(results, query):
    if not results:
        return f"Nothing found for '{query}' in your documents."
    lines = [f"Found in {len(results)} document(s) for '{query}':\n"]
    for i, r in enumerate(results, 1):
        lines.append(f"{i}. {r['name']}")
        lines.append(f"   {r['snippet']}")
        lines.append(f"   Path: {r['path']}\n")
    lines.append("Want me to open any of these? Say the number or filename.")
    return "\n".join(lines)

def build_context_for_query(query, max_chars=8000):
    results, err = search(query, top_n=3)
    if err or not results:
        return None, results
    context_parts = [f"[DOCUMENT SEARCH RESULTS for: {query}]"]
    for r in results:
        context_parts.append(f"\nFile: {r['name']}")
        content = r["content"]
        words = [w.lower() for w in re.findall(r'\b\w+\b', query)]
        snippet_long = _get_snippet(content, words, context=2000)
        context_parts.append(snippet_long)
    full_context = "\n".join(context_parts)
    if len(full_context) > max_chars:
        full_context = full_context[:max_chars]
    return full_context, results

def is_doc_search_request(text):
    keywords = ["notes", "pdf", "ppt", "slides", "document", "file", "lecture",
                 "chapter", "topic", "subject", "search my", "find in",
                 "look in", "check my", "in my files", "in my folder",
                 "handwritten", "scan", "written", "notes on", "notes about",
                 "tell me about", "what does it say", "summarise my",
                 "summarize my", "explain from"]
    tl = text.lower()
    folders = get_folders()
    has_folders = len(folders) > 0
    if not has_folders:
        return False
    return any(k in tl for k in keywords)

def get_index_stats():
    index = _load_index()
    cfg   = _load_config()
    folders = cfg.get("folders", [])
    indexed_at = cfg.get("indexed_at", "never")
    if indexed_at and indexed_at != "never":
        indexed_at = indexed_at[:16].replace("T", " ")
    ext_counts = {}
    for data in index.values():
        ext = data.get("ext", "?")
        ext_counts[ext] = ext_counts.get(ext, 0) + 1
    lines = [
        f"Folders   : {len(folders)}",
        f"Files     : {len(index)}",
        f"Last index: {indexed_at}",
    ]
    for ext, count in sorted(ext_counts.items(), key=lambda x: -x[1]):
        lines.append(f"  {ext}: {count} files")
    return "\n".join(lines)
